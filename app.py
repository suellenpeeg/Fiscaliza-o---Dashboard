import streamlit as st
import pandas as pd
import gspread
from google.oauth2.service_account import Credentials
from pptx import Presentation
from pptx.util import Inches
import plotly.express as px
import tempfile

# ==========================================
# CONFIGURAÇÃO
# ==========================================

st.set_page_config(
    page_title="Dashboard Executivo - Fiscalização 2026",
    layout="wide"
)

ABA_DASHBOARD = "CONTROLE - B. DADOS"

# ==========================================
# CARREGAMENTO SEGURO DOS DADOS
# ==========================================

@st.cache_data(ttl=300)
def load_data():

    try:
        creds = Credentials.from_service_account_info(
            st.secrets["gcp_service_account"],
            scopes=["https://www.googleapis.com/auth/spreadsheets"]
        )

        client = gspread.authorize(creds)

        spreadsheet = client.open_by_url(st.secrets["sheet_url"])
        worksheet = spreadsheet.worksheet(ABA_DASHBOARD)

        data = worksheet.get_all_values()

        if not data or len(data) < 2:
            return pd.DataFrame()

        df = pd.DataFrame(data[1:], columns=data[0])

        # Remove colunas vazias
        df = df.loc[:, df.columns.notna()]
        df = df.loc[:, df.columns != ""]

        # Converte para numérico quando possível
        for col in df.columns:
            df[col] = (
                df[col]
                .astype(str)
                .str.replace(",", ".")
                .str.replace(" ", "")
            )
            df[col] = pd.to_numeric(df[col], errors="coerce")

        # Remove colunas totalmente vazias
        df = df.dropna(axis=1, how="all")

        return df

    except Exception as e:
        st.error(f"Erro ao carregar dados: {e}")
        return pd.DataFrame()


# ==========================================
# EXPORTAÇÃO POWERPOINT
# ==========================================

def gerar_ppt(df_resumo):

    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = "Relatório Executivo - Fiscalização 2026"

    rows, cols = df_resumo.shape

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4)

    table = slide.shapes.add_table(
        rows + 1,
        cols,
        left,
        top,
        width,
        height
    ).table

    # Cabeçalho
    for col in range(cols):
        table.cell(0, col).text = str(df_resumo.columns[col])

    # Dados
    for row in range(rows):
        for col in range(cols):
            table.cell(row + 1, col).text = str(df_resumo.iloc[row, col])

    temp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(temp.name)

    return temp.name


# ==========================================
# APP
# ==========================================

st.title("📊 Dashboard Executivo - Fiscalização 2026")

df = load_data()

if not isinstance(df, pd.DataFrame) or df.empty:
    st.warning("Nenhum dado válido encontrado na planilha.")
    st.stop()

# ==========================================
# IDENTIFICA COLUNAS NUMÉRICAS
# ==========================================

numeric_cols = df.select_dtypes(include="number").columns

if len(numeric_cols) == 0:
    st.warning("Nenhuma coluna numérica válida encontrada para análise.")
    st.dataframe(df)
    st.stop()

# ==========================================
# KPI PRINCIPAL
# ==========================================

total_geral = df[numeric_cols].sum().sum()

col1, col2, col3 = st.columns(3)

col1.metric("Total Geral de Ações", int(total_geral))
col2.metric("Registros (Linhas)", len(df))
col3.metric("Indicadores Numéricos", len(numeric_cols))

st.divider()

# ==========================================
# EVOLUÇÃO DIÁRIA CONSOLIDADA
# ==========================================

df["TOTAL_DIA"] = df[numeric_cols].sum(axis=1)

fig_evolucao = px.line(
    df,
    y="TOTAL_DIA",
    title="Evolução Diária Consolidada",
    markers=True
)

st.plotly_chart(fig_evolucao, use_container_width=True)

st.divider()

# ==========================================
# RANKING POR TIPO
# ==========================================

totais_por_tipo = df[numeric_cols].sum().sort_values(ascending=False)

df_ranking = totais_por_tipo.reset_index()
df_ranking.columns = ["Tipo", "Total"]

fig_ranking = px.bar(
    df_ranking,
    x="Total",
    y="Tipo",
    orientation="h",
    title="Ranking por Tipo de Ação"
)

st.plotly_chart(fig_ranking, use_container_width=True)

st.divider()

# ==========================================
# PARTICIPAÇÃO PERCENTUAL
# ==========================================

fig_pizza = px.pie(
    df_ranking,
    values="Total",
    names="Tipo",
    title="Participação Percentual por Tipo"
)

st.plotly_chart(fig_pizza, use_container_width=True)

st.divider()

# ==========================================
# TABELA COMPLETA
# ==========================================

st.subheader("Base Completa")
st.dataframe(df, use_container_width=True)

st.divider()

# ==========================================
# EXPORTAÇÃO POWERPOINT
# ==========================================

st.subheader("Exportação Executiva")

if st.button("📥 Gerar Relatório em PowerPoint"):

    ppt_file = gerar_ppt(df_ranking)

    with open(ppt_file, "rb") as f:
        st.download_button(
            "Baixar PowerPoint",
            f,
            file_name="relatorio_fiscalizacao_2026.pptx"
        )





